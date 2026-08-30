#!/usr/bin/env python3
"""Build the Solum business-case slide deck (PDF + PPTX) from the same config.json.
Usage: python3 build_deck.py <config.json>
Renders 5 brand slides via Chrome headless -> PDF, writes slideN.png for QA, and a full-bleed PPTX.
ONLY run after the spreadsheet is built and the numbers are agreed, and only if JP asked for slides.
"""
import json, sys, os, base64, subprocess, glob
SKILL_DIR=os.path.dirname(os.path.dirname(os.path.abspath(__file__)))
CHROME="/Applications/Google Chrome.app/Contents/MacOS/Google Chrome"

def usd(x): return f"${x:,.0f}"
def kk(x): return f"${x/1000:,.0f}K"
def mult(x): return f"{x:.1f}×"
def rate_disp(r): return f"${int(r)}" if r==int(r) else f"${r:.2f}"
def vol_disp(p):
    v=p["volume"]; return p.get("volume_label") or (str(int(v)) if v==int(v) else f"{v:g}")

def compute(cfg):
    gt=sum(p["volume"]*p["rate"] for p in cfg["pricing"])
    ot=sum(o["amount"] for o in cfg["one_time"]); ramp=cfg.get("ramp_months",1)
    annual=gt*12; y1c=annual+ot; mode=cfg.get("roi_mode","recurring")
    S=[]
    for i in range(3):
        val=sum(d["rate"]*d["hours"][i] for d in cfg["drivers"]); net=val-gt
        roi=(net/gt) if mode!="year1" else ((val*(12-ramp)-y1c)/y1c)
        S.append(dict(value=val,net=net,roi=roi,ny1=val*(12-ramp)-y1c,
                      pb=ramp+(ot+gt*ramp)/net if net>0 else 0,
                      annual=val*12, yr3=val*(36-ramp)-(gt*36+ot)))
    return dict(gt=gt,ot=ot,ramp=ramp,annual=annual,y1c=y1c,mode=mode,S=S)

CSS="""
@import url('https://fonts.googleapis.com/css2?family=DM+Sans:ital,opsz,wght@0,9..40,400;0,9..40,500;0,9..40,600;0,9..40,700;1,9..40,400&display=swap');
@page{size:13.333in 7.5in;margin:0;}
*{margin:0;padding:0;box-sizing:border-box;-webkit-print-color-adjust:exact;print-color-adjust:exact;}
:root{--navy:#011C40;--blue:#468AF7;--teal:#70D3C6;--ink:#0F1B2D;--muted:#6C7689;--line:#E3E8F0;--alt:#F5F7FA;}
body{font-family:'DM Sans',sans-serif;color:var(--ink);}
.slide{width:13.333in;height:7.5in;position:relative;overflow:hidden;page-break-after:always;background:#fff;display:flex;flex-direction:column;}
.shead{height:.92in;background:var(--navy);display:flex;align-items:center;justify-content:space-between;padding:0 .6in;}
.shead img{height:.34in;} .shead .st{color:#fff;font-weight:700;font-size:17px;}
.shead .se{color:var(--blue);font-weight:700;font-size:9px;letter-spacing:2.5px;text-transform:uppercase;}
.acc{height:5px;background:linear-gradient(90deg,var(--blue),var(--teal));}
.sbody{flex:1;padding:.34in .6in .2in;display:flex;flex-direction:column;}
.sbody.center{justify-content:center;}
.sfoot{height:.36in;background:var(--navy);color:#9DB0CC;display:flex;align-items:center;justify-content:space-between;padding:0 .6in;font-size:8.5px;}
.sfoot b{color:#fff;}
h2.lead{font-size:21px;font-weight:700;color:var(--navy);letter-spacing:-.4px;} h2.lead b{color:var(--blue);}
.sub{font-size:12px;color:var(--muted);margin-top:5px;}
.cover{background:var(--navy);justify-content:center;align-items:center;text-align:center;}
.cover img{height:.66in;margin-bottom:.5in;}
.cover .ey{color:var(--blue);font-weight:700;letter-spacing:5px;font-size:13px;text-transform:uppercase;}
.cover h1{color:#fff;font-size:44px;font-weight:700;letter-spacing:-1px;margin-top:14px;line-height:1.06;}
.cover .csub{color:#AEBBD0;font-size:16px;margin-top:16px;}
.cover .cbar{width:120px;height:5px;background:linear-gradient(90deg,var(--blue),var(--teal));margin-top:24px;border-radius:3px;}
.cover .proof{color:#8595B0;font-size:12px;margin-top:24px;} .cover .cfoot{position:absolute;bottom:.5in;color:#7E8DA8;font-size:11px;}
.hero{display:grid;grid-template-columns:repeat(3,1fr);gap:14px;margin:.2in 0;}
.stat{background:var(--navy);border-radius:11px;padding:18px 20px;color:#fff;position:relative;overflow:hidden;}
.stat:after{content:'';position:absolute;left:0;top:0;bottom:0;width:5px;background:linear-gradient(180deg,var(--blue),var(--teal));}
.stat .k{font-size:10px;letter-spacing:1.3px;text-transform:uppercase;color:#9DB0CC;font-weight:600;}
.stat .v{font-size:40px;font-weight:700;letter-spacing:-1.5px;margin-top:7px;line-height:1;}
.stat .n{font-size:10.5px;color:var(--teal);margin-top:8px;}
.paths{display:grid;grid-template-columns:repeat(3,1fr);gap:14px;margin-top:.06in;}
.path{border:1.5px solid var(--line);border-radius:11px;padding:15px 18px;background:#fff;}
.path.exp{border:2px solid var(--blue);background:#F4F8FF;}
.path .pn{font-size:11px;font-weight:700;letter-spacing:1px;text-transform:uppercase;color:var(--muted);} .path.exp .pn{color:var(--blue);}
.path .pv{font-size:25px;font-weight:700;color:var(--navy);margin-top:6px;}
.path .pr{display:flex;justify-content:space-between;font-size:11px;margin-top:9px;padding-top:8px;border-top:1px solid var(--line);} .path .pr b{color:var(--blue);}
.tag{display:inline-block;background:var(--blue);color:#fff;font-size:8.5px;font-weight:700;letter-spacing:1px;padding:2px 8px;border-radius:10px;margin-left:7px;}
.take{margin-top:.2in;background:#F4F8FF;border-left:4px solid var(--blue);border-radius:8px;padding:13px 18px;font-size:12px;}
.cols2{display:grid;grid-template-columns:1.15fr 1fr;gap:24px;}
.h3{font-size:11px;letter-spacing:1.5px;text-transform:uppercase;color:var(--blue);font-weight:700;margin-bottom:9px;}
table{width:100%;border-collapse:collapse;font-size:12px;}
th{background:var(--navy);color:#fff;font-weight:600;font-size:10px;padding:8px 10px;text-align:left;} th.r,td.r{text-align:right;}
td{padding:7px 10px;border-bottom:1px solid var(--line);} tr:nth-child(even) td{background:var(--alt);}
.tot td{background:var(--navy)!important;color:#fff;font-weight:700;font-size:13px;border:none;}
.box{margin-top:11px;background:#F4F8FF;border:1px solid #D8E4F7;border-radius:8px;padding:11px 14px;font-size:11px;}
.box .bt{font-weight:700;color:#011C40;margin-bottom:6px;letter-spacing:.3px;}
.box .br{display:flex;justify-content:space-between;color:#6C7689;padding:2px 0;} .box .br b{color:#0F1B2D;}
.box .btot{display:flex;justify-content:space-between;border-top:1.5px solid #C9D9F2;margin-top:5px;padding-top:6px;font-weight:700;color:#011C40;font-size:12px;}
.drv{background:var(--alt);border:1px solid var(--line);border-radius:10px;padding:11px 16px;margin-bottom:9px;}
.drv .dt{display:flex;justify-content:space-between;align-items:baseline;} .drv .dn{font-weight:700;font-size:12.5px;color:var(--navy);}
.drv .dr{font-weight:700;font-size:14px;color:var(--blue);} .drv .dc{font-size:10px;color:var(--muted);margin-top:3px;}
.note{font-size:9.5px;color:var(--muted);margin-top:7px;font-style:italic;}
.cmp{margin-top:.22in;background:#fff;border:1px solid var(--line);border-radius:12px;padding:15px 22px;}
.crow{display:flex;align-items:center;gap:12px;margin-bottom:8px;font-size:11px;}
.crow .cl{width:118px;color:var(--muted);font-weight:600;} .crow .ct{flex:1;background:#EEF1F7;border-radius:6px;height:21px;overflow:hidden;}
.crow .cf{height:100%;border-radius:6px;display:flex;align-items:center;padding:0 10px;color:#fff;font-weight:700;font-size:10.5px;}
.gantt{margin-top:.16in;width:100%;border-collapse:collapse;font-size:11px;}
.gantt th{background:var(--navy);color:#fff;font-size:9px;font-weight:600;padding:7px 6px;text-align:center;} .gantt th.l{text-align:left;width:2in;}
.gantt td{padding:0;border:1px solid #EAEEF5;height:.46in;} .gantt td.lab{padding:7px 11px;font-weight:700;color:var(--navy);font-size:11px;background:#fff;}
.gantt td.lab .sm{display:block;font-weight:400;color:var(--muted);font-size:9px;margin-top:2px;}
.bar{height:62%;border-radius:5px;margin:auto;}
.legend{margin-top:13px;font-size:10.5px;color:var(--muted);} .legend span{display:inline-block;width:11px;height:11px;border-radius:3px;vertical-align:middle;margin:0 5px 0 16px;}
.bgrid{display:grid;grid-template-columns:repeat(3,1fr);gap:18px;margin-top:.2in;}
.bcard{background:#fff;border:1px solid var(--line);border-top:5px solid var(--blue);border-radius:12px;padding:24px 22px;}
.bcard .bp{font-size:34px;font-weight:700;color:var(--navy);} .bcard .btt{font-size:14px;font-weight:700;color:var(--blue);margin-top:6px;} .bcard .bd{font-size:11.5px;color:var(--muted);margin-top:8px;line-height:1.4;}
.proofstrip{margin-top:.22in;text-align:center;font-size:11px;color:#6C7689;line-height:1.8;}
.close{margin-top:.26in;background:var(--navy);border-radius:12px;padding:22px 26px;color:#fff;display:flex;justify-content:space-between;align-items:center;}
.close .cl{font-size:18px;font-weight:700;} .close .cl b{color:var(--teal);} .close .cr{font-size:13px;color:#AEBBD0;text-align:right;}
"""

def build_html(cfg,m,L):
    S=m["S"]; gt=m["gt"]; exp=S[1]
    foot=lambda n:f'<div class="sfoot"><span><b>getsolum.com</b> &nbsp;|&nbsp; hello@getsolum.com</span><span>{cfg["client"]} &nbsp;·&nbsp; Confidential &nbsp;·&nbsp; {n}/5</span></div>'
    head=lambda t,e:f'<div class="shead"><img src="{L}"><div style="text-align:right"><div class="se">{e}</div><div class="st">{t}</div></div></div><div class="acc"></div>'
    roi_lbl="Recurring ROI (excl. one-time)" if m["mode"]!="year1" else "Year-1 ROI (incl. setup)"
    roi_note="net return on the platform fee" if m["mode"]!="year1" else "net Yr-1 ÷ Yr-1 investment"
    roih="ROI (recur.)" if m["mode"]!="year1" else "ROI (Yr-1)"
    # cover
    s1=f'''<div class="slide cover"><div><img src="{L}"><div class="ey">Business Case</div>
     <h1>{cfg.get("title","Business Case")}</h1>
     <div class="csub">Prepared for {cfg["client"]} &nbsp;·&nbsp; {cfg.get("date","")}</div><div class="cbar" style="margin:24px auto 0"></div>
     <div class="proof">{cfg.get("proof","")[:120]}</div></div>
     <div class="cfoot">Private &amp; Confidential &nbsp;·&nbsp; getsolum.com</div></div>'''
    # dashboard
    def pathcard(i,nm,exp_=False):
        s=S[i]; cls=" exp" if exp_ else ""; tag='<span class="tag">RECOMMENDED</span>' if exp_ else ""
        return f'<div class="path{cls}"><div class="pn">{nm} {tag}</div><div class="pv">{usd(s["value"])}<span style="font-size:12px;color:#6C7689">/mo value</span></div><div class="pr"><span>{roih}</span><b>{mult(s["roi"])}</b></div><div class="pr"><span>Net / mo</span><b>{usd(s["net"])}</b></div></div>'
    s2=f'''<div class="slide">{head("The Bottom Line","Dashboard")}<div class="sbody">
     <h2 class="lead">A flat <b>{usd(gt)}/month</b> that captures <b>{kk(S[0]["annual"])}–{kk(S[2]["annual"])} a year</b> in value.</h2>
     <div class="sub">Expected case returns {mult(exp["roi"])} on the platform fee; the {usd(m["ot"])} one-time is recovered in about {exp["pb"]:.1f} months. Upside scales to {mult(S[2]["roi"])}.</div>
     <div class="hero">
       <div class="stat"><div class="k">{roi_lbl}</div><div class="v">{mult(exp["roi"])}</div><div class="n">{roi_note}</div></div>
       <div class="stat"><div class="k">Annual Value (expected)</div><div class="v">{kk(exp["annual"])}</div><div class="n">{kk(S[0]["annual"])}–{kk(S[2]["annual"])} across paths</div></div>
       <div class="stat"><div class="k">Payback Period</div><div class="v">~{exp["pb"]:.1f} mo</div><div class="n">incl. setup + ramp</div></div></div>
     <div class="h3">Follow every path</div>
     <div class="paths">{pathcard(0,"Conservative")}{pathcard(1,"Expected",True)}{pathcard(2,"Best case")}</div>
     <div class="take"><b style="color:#011C40">Every path includes the full {usd(m["ot"])} one-time.</b> &nbsp;Over three years, the expected case nets ~{usd(exp["yr3"])} — up to {usd(S[2]["yr3"])}.</div>
     </div>{foot(2)}</div>'''
    # financial impact
    prows="".join(f'<tr><td>{p["service"]}</td><td class="r">{vol_disp(p)}</td><td class="r">{rate_disp(p["rate"])}</td><td class="r">{usd(p["volume"]*p["rate"])}</td></tr>' for p in cfg["pricing"])
    drows="".join(f'<div class="drv"><div class="dt"><div class="dn">{d["name"]}</div><div class="dr">{usd(d["rate"]*d["hours"][0])}–{usd(d["rate"]*d["hours"][2])}</div></div><div class="dc">{d["hours"][0]}–{d["hours"][2]} hrs/mo &times; {rate_disp(d["rate"])}/hr</div></div>' for d in cfg["drivers"])
    sc_rows="".join(f'<tr><td>{n}</td><td class="r">{usd(S[i]["value"])}</td><td class="r">{mult(S[i]["roi"])}</td><td class="r">{usd(S[i]["ny1"])}</td></tr>' for i,n in enumerate(["Conservative","Expected","Best case"]))
    mx=max(s["value"] for s in S)
    def cbar(lbl,val,color,w): return f'<div class="crow"><div class="cl">{lbl}</div><div class="ct"><div class="cf" style="width:{w}%;background:{color}">{usd(val)}</div></div></div>'
    bars=cbar("Platform cost",gt,"#8795AD",max(8,gt/mx*100))+cbar("Conservative",S[0]["value"],"#9CC0F5",S[0]["value"]/mx*100)+cbar("Expected",S[1]["value"],"#468AF7",S[1]["value"]/mx*100)+cbar("Best case",S[2]["value"],"linear-gradient(90deg,#468AF7,#70D3C6)",100)
    ot_rows="".join(f'<div class="br"><span>{o["label"]}</span><b>{usd(o["amount"])}</b></div>' for o in cfg["one_time"])
    s3=f'''<div class="slide">{head("Financial Impact","The Numbers")}<div class="sbody">
     <div class="cols2">
       <div><div class="h3">The Investment — flat monthly</div>
        <table><tr><th>Service</th><th class="r">Vol / mo</th><th class="r">Rate</th><th class="r">Monthly</th></tr>{prows}
        <tr class="tot"><td>Monthly Total</td><td></td><td></td><td class="r">{usd(gt)}</td></tr></table>
        <div class="box"><div class="bt">YEAR-1 COSTS</div><div class="br"><span>Platform — recurring (12 × {usd(gt)})</span><b>{usd(m["annual"])}</b></div>{ot_rows}<div class="btot"><span>Total Year-1 cost</span><span>{usd(m["y1c"])}</span></div></div></div>
       <div><div class="h3">The Value Engine — monthly</div>{drows}
        <table><tr><th>Path</th><th class="r">Value/mo</th><th class="r">{roih}</th><th class="r">Net Yr-1</th></tr>{sc_rows}</table>
        <div class="note">{"ROI = net monthly value ÷ platform fee, excl. one-time; the "+usd(m["ot"])+" setup is recovered via payback." if m["mode"]!="year1" else "ROI = Net Year-1 ÷ total Year-1 investment (incl. setup)."}</div></div>
     </div>
     <div class="cmp"><div class="h3">Monthly value vs the cost — every path</div>{bars}</div>
     </div>{foot(3)}</div>'''
    # timeline
    tl=cfg.get("timeline",[]); wkhdr="".join(f'<th>W{w}</th>' for w in range(9))
    def cells(s,e,kind):
        out=""
        for w in range(9):
            if s<=w<=e: out+=f'<td><div class="bar" style="background:{"#70D3C6" if kind=="validation" else "#468AF7"}"></div></td>'
            else: out+='<td></td>'
        return out
    trows="".join(f'<tr><td class="lab">{p["phase"]}<span class="sm">{p.get("activities","")}</span></td>{cells(p.get("start",0),p.get("end",0),p.get("kind","impl"))}</tr>' for p in tl)
    s4=f'''<div class="slide">{head("Suggested Timeline","Implementation")}<div class="sbody">
     <h2 class="lead">From signature to full automation with a validation window built in.</h2>
     <table class="gantt"><tr><th class="l">Phase</th>{wkhdr}</tr>{trows}</table>
     <div class="legend"><span style="background:#468AF7"></span>Implementation <span style="background:#70D3C6"></span>Validation window</div>
     </div>{foot(4)}</div>''' if tl else ""
    # why it works
    bens=cfg.get("benefits",[])
    bcards="".join(f'<div class="bcard"><div class="bp">{b["stat"]}</div><div class="btt">{b["title"]}</div><div class="bd">{b["desc"]}</div></div>' for b in bens[:3])
    contact=cfg.get("contact","JP Montoya · CEO, Solum Health<br>jp@getsolum.com · (628) 276-2659")
    s5=f'''<div class="slide">{head("Why It Works","The Outcome")}<div class="sbody center">
     <h2 class="lead">{cfg.get("outcome_line","Fewer denials. Faster intake. A freed-up team.")}</h2>
     <div class="bgrid">{bcards}</div>
     <div class="proofstrip">{cfg.get("proof","")}</div>
     <div class="close"><div class="cl">Live in ~4 weeks · <b>{usd(exp["yr3"])}+ net over three years.</b><br><span style="font-size:13px;font-weight:400;color:#AEBBD0">Flat {usd(gt)}/month.</span></div><div class="cr">{contact}</div></div>
     </div>{foot(5)}</div>''' if bens else ""
    return f"<!DOCTYPE html><html><head><meta charset='utf-8'><style>{CSS}</style></head><body>{s1}{s2}{s3}{s4}{s5}</body></html>"

def render(cfg):
    m=compute(cfg); outdir=cfg["output_dir"]; os.makedirs(outdir,exist_ok=True)
    logo=base64.b64encode(open(os.path.join(SKILL_DIR,"assets","solum_logo.png"),"rb").read()).decode()
    L=f"data:image/png;base64,{logo}"
    html=build_html(cfg,m,L); hp=os.path.join(outdir,"_deck.html"); open(hp,"w").write(html)
    pdf=os.path.join(outdir,f'{cfg["client"].replace(" ","_")}_BusinessCase_Deck.pdf')
    subprocess.run([CHROME,"--headless","--disable-gpu","--no-pdf-header-footer",f"--print-to-pdf={pdf}",f"file://{hp}"],check=True,capture_output=True)
    # render PNGs (QA) + hi-res for pptx
    import Quartz
    def png(pg,scale,path):
        url=Quartz.CFURLCreateFromFileSystemRepresentation(None,pdf.encode(),len(pdf.encode()),False)
        doc=Quartz.CGPDFDocumentCreateWithURL(url); page=Quartz.CGPDFDocumentGetPage(doc,pg)
        box=Quartz.CGPDFPageGetBoxRect(page,Quartz.kCGPDFMediaBox); w=int(box.size.width*scale); h=int(box.size.height*scale)
        cs=Quartz.CGColorSpaceCreateDeviceRGB(); ctx=Quartz.CGBitmapContextCreate(None,w,h,8,0,cs,Quartz.kCGImageAlphaPremultipliedLast)
        Quartz.CGContextSetRGBFillColor(ctx,1,1,1,1); Quartz.CGContextFillRect(ctx,Quartz.CGRectMake(0,0,w,h))
        Quartz.CGContextScaleCTM(ctx,scale,scale); Quartz.CGContextDrawPDFPage(ctx,page)
        img=Quartz.CGBitmapContextCreateImage(ctx); d=Quartz.CFURLCreateFromFileSystemRepresentation(None,path.encode(),len(path.encode()),False)
        dest=Quartz.CGImageDestinationCreateWithURL(d,"public.png",1,None); Quartz.CGImageDestinationAddImage(dest,img,None); Quartz.CGImageDestinationFinalize(dest)
    url=Quartz.CFURLCreateFromFileSystemRepresentation(None,pdf.encode(),len(pdf.encode()),False)
    n=Quartz.CGPDFDocumentGetNumberOfPages(Quartz.CGPDFDocumentCreateWithURL(url))
    for pg in range(1,n+1):
        png(pg,1.3,os.path.join(outdir,f"slide{pg}.png"))      # QA
        png(pg,2.0,os.path.join(outdir,f"_hr{pg}.png"))        # pptx
    # pptx
    from pptx import Presentation
    from pptx.util import Inches
    prs=Presentation(); prs.slide_width=Inches(13.333); prs.slide_height=Inches(7.5)
    for pg in range(1,n+1):
        s=prs.slides.add_slide(prs.slide_layouts[6]); s.shapes.add_picture(os.path.join(outdir,f"_hr{pg}.png"),0,0,width=Inches(13.333),height=Inches(7.5))
    pptx=os.path.join(outdir,f'{cfg["client"].replace(" ","_")}_BusinessCase_Deck.pptx'); prs.save(pptx)
    for f in glob.glob(os.path.join(outdir,"_hr*.png"))+[hp]: os.remove(f)
    print("SAVED:",pdf); print("SAVED:",pptx); print(f"QA: review slide1..{n}.png in {outdir}")

if __name__=="__main__":
    if len(sys.argv)<2: print("usage: build_deck.py <config.json>"); sys.exit(1)
    render(json.load(open(sys.argv[1])))
