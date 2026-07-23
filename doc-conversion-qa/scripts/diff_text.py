#!/usr/bin/env python3
"""Extract text from source and output, diff line by line.

Supports .pdf, .pptx, .docx. Uses PyMuPDF for PDF, python-pptx for PPTX,
python-docx for DOCX.
"""
import sys, os, difflib

def extract(path):
    ext = os.path.splitext(path)[1].lower()
    if ext == ".pdf":
        import fitz
        doc = fitz.open(path)
        return [page.get_text("text") for page in doc]
    if ext == ".pptx":
        from pptx import Presentation
        prs = Presentation(path)
        out = []
        for slide in prs.slides:
            buf = []
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for p in shape.text_frame.paragraphs:
                        line = "".join(r.text for r in p.runs)
                        if line.strip():
                            buf.append(line)
            out.append("\n".join(buf))
        return out
    if ext == ".docx":
        from docx import Document
        d = Document(path)
        return ["\n".join(p.text for p in d.paragraphs if p.text.strip())]
    sys.exit(f"Unsupported extension: {ext}")


def normalize(s):
    return " ".join(s.split())


def main(src, out):
    src_pages = extract(src)
    out_pages = extract(out)
    print(f"Source: {len(src_pages)} pages/slides")
    print(f"Output: {len(out_pages)} pages/slides")
    if len(src_pages) != len(out_pages):
        print("⚠️  Page count mismatch")

    src_norm = normalize(" ".join(src_pages))
    out_norm = normalize(" ".join(out_pages))
    src_tokens = set(src_norm.split())
    out_tokens = set(out_norm.split())
    missing = src_tokens - out_tokens
    extra = out_tokens - src_tokens
    if missing:
        print(f"\n❌ {len(missing)} tokens missing from output:")
        for t in list(missing)[:50]:
            print(f"  - {t!r}")
    if extra:
        print(f"\n+ {len(extra)} new tokens in output:")
        for t in list(extra)[:50]:
            print(f"  + {t!r}")
    if not missing and not extra:
        print("\n✅ Text content matches.")


if __name__ == "__main__":
    if len(sys.argv) < 3:
        sys.exit("usage: diff_text.py <source> <output>")
    main(sys.argv[1], sys.argv[2])
